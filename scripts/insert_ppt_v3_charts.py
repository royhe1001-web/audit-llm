"""
insert_ppt_v3_charts.py — 把 3 张图插入到 PPT v3 对应 Slide
"""
import shutil
import os
from pptx import Presentation
from pptx.util import Inches, Emu
from copy import deepcopy
from lxml import etree

PPT_PATH = 'PPT_ 上市公司财务舞弊智能识别与持续审计系统v3.pptx'
CHARTS_DIR = 'output/ppt_charts_v3'

CHARTS = {
    6: 'feature_importance_top10.png',   # Slide 6: 特征重要性 TOP 10
    9: 'f1_evolution.png',                # Slide 9: F1 演进柱状图
    11: 'ablation_study.png',             # Slide 11: 消融实验
}

def add_picture_to_slide(slide, image_path, position='center'):
    """在 Slide 中部添加图片"""
    pic = slide.shapes.add_picture(
        image_path,
        Inches(2.5), Inches(2.5),   # 左上角位置
        Inches(8.0), Inches(4.0),   # 大小
    )
    return pic


def main():
    ppt = Presentation(PPT_PATH)
    print(f'原 PPT: {len(ppt.slides)} 张 Slide')

    for slide_idx, chart_name in CHARTS.items():
        chart_path = os.path.join(CHARTS_DIR, chart_name)
        if not os.path.exists(chart_path):
            print(f'  ⚠️ 找不到图表: {chart_path}')
            continue

        slide = ppt.slides[slide_idx - 1]  # 0-based
        # 在 Slide 中部添加图片
        add_picture_to_slide(slide, chart_path)
        print(f'  ✅ Slide {slide_idx}: 插入 {chart_name}')

    # 保存(另存为 v3.1)
    out_path = PPT_PATH.replace('v3.pptx', 'v3.1.pptx')
    ppt.save(out_path)
    print(f'\n✅ 已保存: {out_path}')


if __name__ == '__main__':
    main()
